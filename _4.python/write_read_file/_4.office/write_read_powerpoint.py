"""
讀寫 Powerpoint 檔案, 使用 pptx
pip install python-pptx
pip install python-pptx

"""

import sys

print("------------------------------------------------------------")  # 60個

import pptx
from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.shapes import MSO_SHAPE

#建立新的簡報物件
pres=Presentation()
#增加第一張投影片(套用母片樣板0-標題頁)
slide=pres.slides.add_slide(pres.slide_layouts[0])
#設定第一張投影片的背景顏色(不過建議先設定好佈景主題，直接開啟簡報檔新增)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb=RGBColor(0,255,255)
#設定第一張投影片的標題
slide.shapes.title.text="標題頁-標題"
#設定第一張投影片的副標題
slide.placeholders[1].text_frame.text="標題頁-副標題"
#增加第二張投影片(套用母片樣板1-標題及物件頁)
slide=pres.slides.add_slide(pres.slide_layouts[1])
#設定第二張投影片的標題
slide.shapes.title.text="標題及物件頁-標題"
#設定第一張投影片物件的內容
textframe=slide.shapes.placeholders[1].text_frame
textframe.text='第一大點'
paragraph=textframe.add_paragraph()
paragraph.text='一之一'
paragraph.level=1
paragraph=textframe.add_paragraph()
paragraph.text='一之二'
paragraph.level=1
paragraph=textframe.add_paragraph()
paragraph.text='第二大點'
paragraph.level=0
paragraph=textframe.add_paragraph()
paragraph.text='二之一'
paragraph.level=1
paragraph=textframe.add_paragraph()
paragraph.text='二之二'
paragraph.level=1
slide.shapes.add_picture('data/Google-Colab-Guide-1024x683.jpg', Inches(3), Inches(3), height=Inches(4.5))
#增加第三張投影片(套用母片樣板2-區段標題頁)
slide=pres.slides.add_slide(pres.slide_layouts[2])
#設定第三張投影片的標題
slide.shapes.title.text="區段標題頁-標題"
#設定第三張投影片的文字區段內容
slide.shapes.placeholders[1].text_frame.text="區段標題頁-區段文字"
#增加第四張投影片(套用母片樣板3-兩項物件頁)
slide=pres.slides.add_slide(pres.slide_layouts[3])
#設定第四張投影片的標題
slide.shapes.title.text="兩項物件頁-標題"
#設定第四張投影片的文字物件內容
slide.shapes.placeholders[1].text_frame.text="兩項物件頁-區段1文字"
slide.shapes.placeholders[2].text_frame.text="兩項物件頁-區段2文字"
#增加第五張投影片(套用母片樣板4-比對頁)
slide=pres.slides.add_slide(pres.slide_layouts[4])
#設定第五張投影片的標題
slide.shapes.title.text="比對頁-標題"
#設定第五張投影片的文字物件內容
slide.shapes.placeholders[1].text_frame.text="比對頁-區段1文字"
slide.shapes.placeholders[2].text_frame.text="比對頁-區段2文字"
slide.shapes.placeholders[3].text_frame.text="比對頁-區段3文字"
slide.shapes.placeholders[4].text_frame.text="比對頁-區段4文字"
#增加第六張投影片(套用母片樣板5-只有標題頁)
slide=pres.slides.add_slide(pres.slide_layouts[5])
#設定第六張投影片的標題
slide.shapes.title.text="只有標題頁-標題"
#增加第七張投影片(套用母片樣板6-空白頁)
slide=pres.slides.add_slide(pres.slide_layouts[6])
#隨便加入物件
#表格: 列，欄，左，上，寬，高
shape=slide.shapes.add_table(3,4,Cm(1),Cm(1),Cm(6),Cm(6))
shape.table.cell(0,0).text="(0,0)"
shape.table.cell(1,1).text="(1,1)"
shape.table.cell(2,3).text="(2,3)"
#圖案: MSO_SHAPE ID，左，上，寬，高
slide.shapes.add_shape(MSO_SHAPE.CUBE,Cm(1),Cm(8),Cm(6),Cm(6))
slide.shapes.add_shape(MSO_SHAPE.SMILEY_FACE,Cm(10),Cm(1),Cm(6),Cm(6))
#圖片: 檔名，左，上(，width=寬，height=高)
slide.shapes.add_picture('data/Google-Colab-Guide-1024x683.jpg',Inches(3),Inches(3),height=Inches(4.5))
#增加第八張投影片(套用母片樣板7-含標題的內容頁)
slide=pres.slides.add_slide(pres.slide_layouts[7])
slide.shapes.title.text="含標題的內容頁-標題"
#設定第八張投影片的文字物件內容
slide.shapes.placeholders[1].text_frame.text="含標題的內容頁-區段1文字"
slide.shapes.placeholders[2].text_frame.text="含標題的內容頁-區段2文字"
#增加第九張投影片(套用母片樣板8-含標題的圖片頁)
slide=pres.slides.add_slide(pres.slide_layouts[8])
#設定第九張投影片的標題
slide.shapes.title.text="含標題的圖片頁-標題"
#設定第九張投影片的圖片
slide.shapes.placeholders[1].insert_picture('data/Google-Colab-Guide-1024x683.jpg')
slide.shapes.placeholders[2].text_frame.text="含標題的圖片頁-區段1文字"
#增加第十張投影片(套用母片樣板9-標題及直排文字頁)
slide=pres.slides.add_slide(pres.slide_layouts[9])
#設定第十張投影片的標題
slide.shapes.title.text="標題及直排文字頁-標題"
#設定第十張投影片文字物件的內容
textframe=slide.shapes.placeholders[1].text_frame
textframe.text='第一大點'
paragraph=textframe.add_paragraph()
paragraph.text='一之一'
paragraph.level=1
paragraph=textframe.add_paragraph()
paragraph.text='一之二'
paragraph.level=1
paragraph=textframe.add_paragraph()
paragraph.text='第二大點'
paragraph.level=0
paragraph=textframe.add_paragraph()
paragraph.text='二之一'
paragraph.level=1
paragraph=textframe.add_paragraph()
paragraph.text='二之二'
paragraph.level=1
#增加第十一張投影片(套用母片樣板10-直排標題及文字頁)
slide=pres.slides.add_slide(pres.slide_layouts[10])
#設定第十一張投影片的標題
slide.shapes.title.text="直排標題及文字頁-標題"
#設定第十一張投影片文字物件的內容
textframe=slide.shapes.placeholders[1].text_frame
textframe.text='第一大點'
paragraph=textframe.add_paragraph()
paragraph.text='一之一'
paragraph.level=1
paragraph=textframe.add_paragraph()
paragraph.text='一之二'
paragraph.level=1
paragraph=textframe.add_paragraph()
paragraph.text='第二大點'
paragraph.level=0
paragraph=textframe.add_paragraph()
paragraph.text='二之一'
paragraph.level=1
paragraph=textframe.add_paragraph()
paragraph.text='二之二'
paragraph.level=1
#儲存簡報檔
pres.save("tmp_powerpoint.pptx")

print("------------------------------------------------------------")  # 60個




print("------------------------------------------------------------")  # 60個


print("------------------------------------------------------------")  # 60個





print("------------------------------------------------------------")  # 60個


print("------------------------------------------------------------")  # 60個



print("------------------------------------------------------------")  # 60個
print("作業完成")
print("------------------------------------------------------------")  # 60個
