"""
讀寫 Word 檔案, 使用 docx
pip install python-docx

"""
import sys
import docx

print("------------------------------------------------------------")  # 60個
print("------------------------------------------------------------")  # 60個

filename_r = "data/python_docx1.docx"
print("讀取 docx, 檔案 : " + filename_r)

doc = docx.Document(filename_r)

for i, paragraph in enumerate(doc.paragraphs):
    print(paragraph.text)

for p in doc.paragraphs:
    print(p.text)

print("------------------------------------------------------------")  # 60個
print("------------------------------------------------------------")  # 60個

# 讀取excel檔, 製作word檔

import openpyxl
from docxtpl import DocxTemplate
import datetime

filename_r1 = "data/python_ReadWrite_EXCEL6_student_data2.xlsx"
print("讀取 xlsx, 檔案 : " + filename_r1)

workbook = openpyxl.load_workbook(filename_r1)
sheet = workbook.active
list_values = list(sheet.values)
print(list_values)

filename_r2 = "data/python_ReadWrite_WORD1_certificate.docx"
print("讀取 docx, 檔案 : " + filename_r2)

# Generate docs
doc = DocxTemplate(filename_r2)
for value_tuple in list_values[1:]:
    doc.render(
        {
            "name": value_tuple[0],
            "course": value_tuple[1],
            "date": value_tuple[2],
            "instructor": value_tuple[3],
        }
    )

    doc_name = "tmp_certificate" + value_tuple[0] + value_tuple[1] + ".docx"
    print("寫入 docx, 檔案 : " + doc_name)
    doc.save(doc_name)

print("------------------------------------------------------------")  # 60個
print("------------------------------------------------------------")  # 60個

"""
讀取word檔, 製作word檔
invoice : 發票，發貨單
"""

import datetime
from docxtpl import DocxTemplate

filename_r = "data/python_ReadWrite_WORD2_invoice_template.docx"
print("讀取 docx, 檔案 : " + filename_r)

doc = DocxTemplate(filename_r)

invoice_list = [[2, "pen", 0.5, 1], [1, "paper pack", 5, 5], [2, "notebook", 2, 4]]

"""
doc.render({"name":"john", 
            "phone":"555-55555",
            "invoice_list": invoice_list,
            "subtotal":10,
            "salestax":"10%",
            "total":9})
"""
name = "David Wang"
phone = "0912-345678"
subtotal = sum(item[3] for item in invoice_list)
salestax = 0.1
total = subtotal * (1 - salestax)

doc.render(
    {
        "name": name,
        "phone": phone,
        "invoice_list": invoice_list,
        "subtotal": subtotal,
        "salestax": str(salestax * 100) + "%",
        "total": total,
    }
)

filename_w = (
    "tmp_new_invoice"
    + name
    + datetime.datetime.now().strftime("%Y-%m-%d-%H%M%S")
    + ".docx"
)
doc.save(filename_w)

print("------------------------------------------------------------")  # 60個
print("------------------------------------------------------------")  # 60個

# 讀取Word文件

filename_r = "./data/用函數還是用復雜的表達式.docx"
print("讀取 docx, 檔案 : " + filename_r)

doc = docx.Document(filename_r)
print(len(doc.paragraphs))
print(doc.paragraphs[0].text)
# print(doc.paragraphs[1].runs[0].text)

content = []
for para in doc.paragraphs:
    content.append(para.text)
print("".join(content))

print("------------------------------------------------------------")  # 60個
print("------------------------------------------------------------")  # 60個

from docx.shared import Cm
from docx.shared import Inches

# 開新的Word檔案
doc = docx.Document()
# 標題 段落樣式
doc.add_heading("標題 段落樣式", 0)
# 內文 段落樣式
p = doc.add_paragraph("內文 段落樣式包含一些 ")
p.add_run("粗體").bold = True
p.add_run(" 還有一些 ")
p.add_run("斜體字.").italic = True
# 標題1 段落樣式
doc.add_heading("標題1 段落樣式", level=1)
# 鮮明引文 段落樣式
doc.add_paragraph("鮮明引文 段落樣式", style="Intense Quote")
# 項目符號 段落樣式
doc.add_paragraph("項目符號 段落樣式,1", style="List Bullet")
doc.add_paragraph("項目符號 段落樣式,2", style="List Bullet")
doc.add_paragraph("清單號碼 段落樣式,1", style="List Number")
doc.add_paragraph("清單號碼 段落樣式,2", style="List Number")
# 插入圖片，指定寬度5cm
doc.add_picture("data/Google-Colab-Guide-1024x683.jpg", width=Cm(5))
# 插入表格，先增加1列標題列
table = doc.add_table(rows=1, cols=3, style="Table Grid")
# 表格標題列
header_cells = table.rows[0].cells
header_cells[0].text = "班級"
header_cells[1].text = "座號"
header_cells[2].text = "姓名"
# 表格內容第1列
row_cells = table.add_row().cells
row_cells[0].text = "101"
row_cells[1].text = "01"
row_cells[2].text = "劉的華"
# 表格內容第2列
row_cells = table.add_row().cells
row_cells[0].text = "101"
row_cells[1].text = "02"
row_cells[2].text = "郭付錢"
# 換頁符號
doc.add_page_break()
# 儲存檔案
doc.save("tmp_word1.docx")

print("------------------------------------------------------------")  # 60個
print("------------------------------------------------------------")  # 60個


print("------------------------------------------------------------")  # 60個
print("------------------------------------------------------------")  # 60個

"""
讀寫 Powerpoint 檔案, 使用 pptx
pip install python-pptx
pip install python-pptx

"""
print("------------------------------------------------------------")  # 60個
import pptx
from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 建立新的簡報物件
pres = Presentation()
# 增加第一張投影片(套用母片樣板0-標題頁)
slide = pres.slides.add_slide(pres.slide_layouts[0])
# 設定第一張投影片的背景顏色(不過建議先設定好佈景主題，直接開啟簡報檔新增)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0, 255, 255)
# 設定第一張投影片的標題
slide.shapes.title.text = "標題頁-標題"
# 設定第一張投影片的副標題
slide.placeholders[1].text_frame.text = "標題頁-副標題"
# 增加第二張投影片(套用母片樣板1-標題及物件頁)
slide = pres.slides.add_slide(pres.slide_layouts[1])
# 設定第二張投影片的標題
slide.shapes.title.text = "標題及物件頁-標題"
# 設定第一張投影片物件的內容
textframe = slide.shapes.placeholders[1].text_frame
textframe.text = "第一大點"
paragraph = textframe.add_paragraph()
paragraph.text = "一之一"
paragraph.level = 1
paragraph = textframe.add_paragraph()
paragraph.text = "一之二"
paragraph.level = 1
paragraph = textframe.add_paragraph()
paragraph.text = "第二大點"
paragraph.level = 0
paragraph = textframe.add_paragraph()
paragraph.text = "二之一"
paragraph.level = 1
paragraph = textframe.add_paragraph()
paragraph.text = "二之二"
paragraph.level = 1
slide.shapes.add_picture(
    "data/Google-Colab-Guide-1024x683.jpg", Inches(3), Inches(3), height=Inches(4.5)
)
# 增加第三張投影片(套用母片樣板2-區段標題頁)
slide = pres.slides.add_slide(pres.slide_layouts[2])
# 設定第三張投影片的標題
slide.shapes.title.text = "區段標題頁-標題"
# 設定第三張投影片的文字區段內容
slide.shapes.placeholders[1].text_frame.text = "區段標題頁-區段文字"
# 增加第四張投影片(套用母片樣板3-兩項物件頁)
slide = pres.slides.add_slide(pres.slide_layouts[3])
# 設定第四張投影片的標題
slide.shapes.title.text = "兩項物件頁-標題"
# 設定第四張投影片的文字物件內容
slide.shapes.placeholders[1].text_frame.text = "兩項物件頁-區段1文字"
slide.shapes.placeholders[2].text_frame.text = "兩項物件頁-區段2文字"
# 增加第五張投影片(套用母片樣板4-比對頁)
slide = pres.slides.add_slide(pres.slide_layouts[4])
# 設定第五張投影片的標題
slide.shapes.title.text = "比對頁-標題"
# 設定第五張投影片的文字物件內容
slide.shapes.placeholders[1].text_frame.text = "比對頁-區段1文字"
slide.shapes.placeholders[2].text_frame.text = "比對頁-區段2文字"
slide.shapes.placeholders[3].text_frame.text = "比對頁-區段3文字"
slide.shapes.placeholders[4].text_frame.text = "比對頁-區段4文字"
# 增加第六張投影片(套用母片樣板5-只有標題頁)
slide = pres.slides.add_slide(pres.slide_layouts[5])
# 設定第六張投影片的標題
slide.shapes.title.text = "只有標題頁-標題"
# 增加第七張投影片(套用母片樣板6-空白頁)
slide = pres.slides.add_slide(pres.slide_layouts[6])
# 隨便加入物件
# 表格: 列，欄，左，上，寬，高
shape = slide.shapes.add_table(3, 4, Cm(1), Cm(1), Cm(6), Cm(6))
shape.table.cell(0, 0).text = "(0,0)"
shape.table.cell(1, 1).text = "(1,1)"
shape.table.cell(2, 3).text = "(2,3)"
# 圖案: MSO_SHAPE ID，左，上，寬，高
slide.shapes.add_shape(MSO_SHAPE.CUBE, Cm(1), Cm(8), Cm(6), Cm(6))
slide.shapes.add_shape(MSO_SHAPE.SMILEY_FACE, Cm(10), Cm(1), Cm(6), Cm(6))
# 圖片: 檔名，左，上(，width=寬，height=高)
slide.shapes.add_picture(
    "data/Google-Colab-Guide-1024x683.jpg", Inches(3), Inches(3), height=Inches(4.5)
)
# 增加第八張投影片(套用母片樣板7-含標題的內容頁)
slide = pres.slides.add_slide(pres.slide_layouts[7])
slide.shapes.title.text = "含標題的內容頁-標題"
# 設定第八張投影片的文字物件內容
slide.shapes.placeholders[1].text_frame.text = "含標題的內容頁-區段1文字"
slide.shapes.placeholders[2].text_frame.text = "含標題的內容頁-區段2文字"
# 增加第九張投影片(套用母片樣板8-含標題的圖片頁)
slide = pres.slides.add_slide(pres.slide_layouts[8])
# 設定第九張投影片的標題
slide.shapes.title.text = "含標題的圖片頁-標題"
# 設定第九張投影片的圖片
slide.shapes.placeholders[1].insert_picture("data/Google-Colab-Guide-1024x683.jpg")
slide.shapes.placeholders[2].text_frame.text = "含標題的圖片頁-區段1文字"
# 增加第十張投影片(套用母片樣板9-標題及直排文字頁)
slide = pres.slides.add_slide(pres.slide_layouts[9])
# 設定第十張投影片的標題
slide.shapes.title.text = "標題及直排文字頁-標題"
# 設定第十張投影片文字物件的內容
textframe = slide.shapes.placeholders[1].text_frame
textframe.text = "第一大點"
paragraph = textframe.add_paragraph()
paragraph.text = "一之一"
paragraph.level = 1
paragraph = textframe.add_paragraph()
paragraph.text = "一之二"
paragraph.level = 1
paragraph = textframe.add_paragraph()
paragraph.text = "第二大點"
paragraph.level = 0
paragraph = textframe.add_paragraph()
paragraph.text = "二之一"
paragraph.level = 1
paragraph = textframe.add_paragraph()
paragraph.text = "二之二"
paragraph.level = 1
# 增加第十一張投影片(套用母片樣板10-直排標題及文字頁)
slide = pres.slides.add_slide(pres.slide_layouts[10])
# 設定第十一張投影片的標題
slide.shapes.title.text = "直排標題及文字頁-標題"
# 設定第十一張投影片文字物件的內容
textframe = slide.shapes.placeholders[1].text_frame
textframe.text = "第一大點"
paragraph = textframe.add_paragraph()
paragraph.text = "一之一"
paragraph.level = 1
paragraph = textframe.add_paragraph()
paragraph.text = "一之二"
paragraph.level = 1
paragraph = textframe.add_paragraph()
paragraph.text = "第二大點"
paragraph.level = 0
paragraph = textframe.add_paragraph()
paragraph.text = "二之一"
paragraph.level = 1
paragraph = textframe.add_paragraph()
paragraph.text = "二之二"
paragraph.level = 1

# 儲存簡報檔
filename_w = "tmp_powerpoint.pptx"
pres.save(filename_w)
print("寫入 pptx, 檔案 : " + filename_w)

print("------------------------------------------------------------")  # 60個
print("------------------------------------------------------------")  # 60個


print("------------------------------------------------------------")  # 60個
print("------------------------------------------------------------")  # 60個


print("------------------------------------------------------------")  # 60個
print("作業完成")
print("------------------------------------------------------------")  # 60個
sys.exit()


print("------------------------------------------------------------")  # 60個


print("------------------------------------------------------------")  # 60個


print("------------------------------------------------------------")  # 60個
